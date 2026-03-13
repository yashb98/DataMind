"""
MCP Report Generator — python-pptx PPTX generation.
Day 11: Phase 2 — Structured PowerPoint output from report sections.

Protocols: None (internal utility called by main.py MCP tool)
SOLID: SRP (PPTX generation only), OCP (extend PPTXGenerator for new slide layouts)
Benchmark: tests/benchmarks/bench_report.py — target < 2s for 10-slide deck
"""

from __future__ import annotations

import asyncio
import io
import time
from typing import Any

import structlog
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

log = structlog.get_logger(__name__)

# ── Brand colours ─────────────────────────────────────────────────────────────
_NAVY = RGBColor(0x1A, 0x36, 0x5D)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xF7, 0xFA, 0xFC)
_DARK_GRAY = RGBColor(0x71, 0x80, 0x96)
_BLUE_ACCENT = RGBColor(0x4A, 0x90, 0xD9)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
_SLIDE_WIDTH = Inches(13.33)
_SLIDE_HEIGHT = Inches(7.5)

# Maximum rows in a table on a slide (header + data rows)
_MAX_TABLE_ROWS = 8
# Maximum characters of section content shown on a slide body
_MAX_BODY_CHARS = 600


def _set_text_color(run: Any, color: RGBColor) -> None:
    """Set font colour on a text run.

    Args:
        run: pptx Run object.
        color: RGBColor to apply.
    """
    run.font.color.rgb = color


def _add_title_slide(prs: Presentation, title: str, generated_at: str) -> None:
    """Add the opening title slide with report name and generation timestamp.

    Args:
        prs: Active Presentation instance.
        title: Report title string.
        generated_at: ISO-8601 UTC timestamp for the subtitle.
    """
    slide_layout = prs.slide_layouts[0]  # "Title Slide" layout
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title

    # Subtitle placeholder (index 1)
    subtitle_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if subtitle_ph is not None:
        subtitle_ph.text = f"DataMind Enterprise Report\nGenerated: {generated_at}"


def _add_content_slide(
    prs: Presentation,
    heading: str,
    content: str,
    data: list[dict[str, Any]] | None,
) -> None:
    """Add a content slide for one report section.

    If `data` is provided and non-empty, a formatted table is added below the
    content text. Tables are capped at `_MAX_TABLE_ROWS` rows and rendered
    with the DataMind navy header style.

    Args:
        prs: Active Presentation instance.
        heading: Section heading rendered as the slide title.
        content: Plain-text body (Markdown stripped for slide display).
        data: Optional list of row dicts to render as a table.
    """
    content_layout = prs.slide_layouts[1]  # "Title and Content" layout
    slide = prs.slides.add_slide(content_layout)

    # Slide title
    slide.shapes.title.text = heading

    # Body text — truncate to avoid overflow
    body_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body_ph is not None:
        # Strip common Markdown syntax for cleaner slide display
        clean_content = _strip_markdown(content)
        tf = body_ph.text_frame
        tf.text = clean_content[:_MAX_BODY_CHARS]
        tf.word_wrap = True
        # Set body font size
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)

    # Optional data table — placed in lower half of slide
    if data and len(data) > 0:
        _add_data_table(slide, data)


def _add_data_table(slide: Any, data: list[dict[str, Any]]) -> None:
    """Add a formatted data table in the lower region of a content slide.

    Args:
        slide: pptx Slide object to add the table to.
        data: List of row dicts; column headers derived from first row keys.
    """
    cols = list(data[0].keys())
    row_count = min(len(data) + 1, _MAX_TABLE_ROWS)  # +1 for header

    left = Inches(0.5)
    top = Inches(4.2)
    width = Inches(12.33)
    height = Inches(2.8)

    table = slide.shapes.add_table(row_count, len(cols), left, top, width, height).table

    # Header row — navy background, white bold text
    for ci, col_name in enumerate(cols):
        cell = table.cell(0, ci)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = _WHITE
                run.font.bold = True
                run.font.size = Pt(10)

    # Data rows
    for ri, row in enumerate(data[: _MAX_TABLE_ROWS - 1], start=1):
        for ci, col_name in enumerate(cols):
            cell = table.cell(ri, ci)
            cell.text = str(row.get(col_name, ""))
            # Alternate row shading
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _LIGHT_GRAY
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)


def _add_provenance_slide(prs: Presentation, report_id: str, merkle_root: str) -> None:
    """Add a provenance certificate as the final slide.

    Args:
        prs: Active Presentation instance.
        report_id: Unique report identifier.
        merkle_root: SHA-256 Merkle root hash of the report content.
    """
    blank_layout = prs.slide_layouts[5]  # "Blank" layout
    slide = prs.slides.add_slide(blank_layout)

    # Title text box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Provenance Certificate"
    for para in title_tf.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(20)
            run.font.color.rgb = _NAVY

    # Divider line (thin rectangle)
    slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5),
        Inches(1.1),
        Inches(12.33),
        Inches(0.02),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _BLUE_ACCENT
    slide.shapes[-1].line.fill.background()

    # Content text box
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    lines = [
        ("Platform:", "DataMind Enterprise v2 — Verifiable AI Analytics"),
        ("Report ID:", report_id),
        ("Merkle Root (SHA-256):", merkle_root),
        ("Protocol:", "MCP + A2A — AAIF / Linux Foundation Standard"),
        (
            "Integrity:",
            "This document's content is cryptographically sealed by the Merkle root "
            "above. Any tampering will invalidate the hash. Anchor verification "
            "available via IPFS.",
        ),
    ]

    for i, (label, value) in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.space_before = Pt(6)
        run_label = para.add_run()
        run_label.text = f"{label}  "
        run_label.font.bold = True
        run_label.font.size = Pt(11)
        run_label.font.color.rgb = _NAVY
        run_value = para.add_run()
        run_value.text = value
        run_value.font.size = Pt(10)
        run_value.font.color.rgb = _DARK_GRAY


def _strip_markdown(text: str) -> str:
    """Remove common Markdown syntax for plain-text slide display.

    Removes headings (#), bold/italic markers (*, _), and code fences.

    Args:
        text: Markdown-formatted string.

    Returns:
        Cleaned plain text suitable for PowerPoint body copy.
    """
    import re

    # Strip headings
    text = re.sub(r"^#+\s+", "", text, flags=re.MULTILINE)
    # Strip bold/italic
    text = re.sub(r"\*{1,3}(.+?)\*{1,3}", r"\1", text)
    text = re.sub(r"_{1,3}(.+?)_{1,3}", r"\1", text)
    # Strip inline code
    text = re.sub(r"`(.+?)`", r"\1", text)
    # Strip fenced code blocks
    text = re.sub(r"```[\s\S]*?```", "[code block]", text)
    # Strip horizontal rules
    text = re.sub(r"^---+$", "", text, flags=re.MULTILINE)
    return text.strip()


def _build_presentation_sync(
    report_id: str,
    title: str,
    sections: list[dict[str, Any]],
    merkle_root: str,
    generated_at: str,
) -> tuple[bytes, int]:
    """Build PPTX bytes synchronously (called inside thread-pool executor).

    Args:
        report_id: Unique report identifier.
        title: Report title for title slide.
        sections: List of section dicts (heading, content, data).
        merkle_root: SHA-256 Merkle root for provenance slide.
        generated_at: ISO-8601 UTC timestamp.

    Returns:
        Tuple of (pptx_bytes, slide_count).
    """
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    _add_title_slide(prs, title, generated_at)

    for section in sections:
        _add_content_slide(
            prs,
            heading=section.get("heading", ""),
            content=section.get("content", ""),
            data=section.get("data"),
        )

    _add_provenance_slide(prs, report_id, merkle_root)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return buf.getvalue(), len(prs.slides)


async def generate_pptx(
    report_id: str,
    title: str,
    sections: list[dict[str, Any]],
    merkle_root: str,
    generated_at: str,
) -> tuple[bytes, int]:
    """Generate a PPTX presentation from structured report sections.

    Creates a widescreen (13.33" x 7.5") PowerPoint presentation with:
    - An opening title slide.
    - One content slide per section (with optional data table).
    - A final provenance certificate slide.

    python-pptx is CPU-bound; the actual build runs in a thread-pool executor
    to keep the asyncio event loop unblocked.

    Args:
        report_id: Unique report identifier embedded in the provenance slide.
        title: Report title for the cover slide.
        sections: List of section dicts with keys: heading, content, data.
        merkle_root: SHA-256 Merkle root to display in provenance slide.
        generated_at: ISO-8601 UTC timestamp shown in the title subtitle.

    Returns:
        Tuple of (pptx_bytes, slide_count). slide_count includes title and
        provenance slides.

    Raises:
        RuntimeError: If python-pptx fails to build the presentation.
    """
    start = time.perf_counter()
    loop = asyncio.get_event_loop()

    try:
        pptx_bytes, slide_count = await loop.run_in_executor(
            None,
            _build_presentation_sync,
            report_id,
            title,
            sections,
            merkle_root,
            generated_at,
        )
    except Exception as exc:
        log.error("pptx.render.failed", error=str(exc), report_id=report_id)
        raise RuntimeError(f"python-pptx PPTX render failed: {exc}") from exc

    elapsed_ms = (time.perf_counter() - start) * 1000
    log.info(
        "pptx.generated",
        report_id=report_id,
        size_bytes=len(pptx_bytes),
        slide_count=slide_count,
        elapsed_ms=round(elapsed_ms, 2),
    )

    return pptx_bytes, slide_count
